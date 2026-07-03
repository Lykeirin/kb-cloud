#!/usr/bin/env python3
"""
知识库文件监控守护进程

监控 ~/kb-inbox/ 目录，新文件放入后自动：
1. 解析文件 → 提取元数据 → 打标签
2. 录入 PostgreSQL 知识库
3. 移动到 ~/kb-archive/ 归档

支持格式：PDF, TXT, Markdown, DOCX, EPUB, HTML, PPTX, 图片(OCR), 压缩包(ZIP/TAR.GZ/TAR.BZ2/TAR.XZ)

用法：
  # 前台运行（调试）
  python kb_watcher.py

  # 后台运行
  nohup python kb_watcher.py > ~/kb-watcher.log 2>&1 &
"""

import sys
import os
import re
import time
import shutil
import signal
import logging
from pathlib import Path
from datetime import datetime

# ─── 数据库配置（硬编码以避免子进程环境变量丢失）───
os.environ.setdefault("KB_DB_HOST", "localhost")
os.environ.setdefault("KB_DB_PORT", "5433")
os.environ.setdefault("KB_DB_NAME", "knowledge_base")
os.environ.setdefault("KB_DB_USER", "kb_admin")
os.environ.setdefault("KB_DB_PASSWORD", "change_me_please")

# 添加模块路径（确保能找到同目录下的 pdf_extractor, kb_core 等模块）
MODULE_DIR = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, MODULE_DIR)

import fitz  # pymupdf

from watchdog.observers import Observer
from watchdog.events import FileSystemEventHandler

from pdf_extractor import (
    extract_title,
    extract_author,
    extract_summary,
    extract_keywords,
    extract_year,
    auto_tag_enhanced,
)
from kb_core import KnowledgeBase


# ─── 配置 ───
WATCH_DIR = Path(os.getenv("KB_INBOX_DIR", str(Path.home() / "kb-inbox")))
ARCHIVE_DIR = Path(os.getenv("KB_ARCHIVE_DIR", str(Path.home() / "kb-archive")))
LOG_FILE = Path(os.getenv("KB_LOG_FILE", str(Path.home() / "kb-watcher.log")))

SUPPORTED_EXTENSIONS = {".pdf", ".txt", ".md", ".markdown", ".docx", ".epub", ".html", ".htm", ".pptx", ".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".webp"}

# 压缩包格式（用文件名后缀匹配，因为 .tar.gz 的 suffix 只返回 .gz）
ARCHIVE_EXTENSIONS = {".zip", ".tar.gz", ".tgz", ".tar.bz2", ".tar.xz"}


def is_archive(filepath: Path) -> bool:
    """检查文件是否为支持的压缩包格式"""
    name = filepath.name.lower()
    return any(name.endswith(ext) for ext in ARCHIVE_EXTENSIONS)

# ─── 日志 ───
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    handlers=[
        logging.FileHandler(LOG_FILE),
        logging.StreamHandler(sys.stdout),
    ],
)
log = logging.getLogger("kb_watcher")


# ─── 文档类型推测 ───
def guess_doc_type(text: str, filename: str, domain: str = "") -> str:
    """推测文档类型"""
    head = text[:2000]
    filename_lower = filename.lower()

    # 写作域
    if domain == "writing" or any(
        k in filename_lower for k in ["chapter", "章节", "人物", "大纲", "设定", "序章", "楔子", "番外"]
    ):
        if any(k in head for k in ["人物", "性格", "外貌", "背景", "角色设定"]):
            return "character"
        if any(k in head for k in ["大纲", "主线", "支线", "剧情走向"]):
            return "outline"
        if any(k in head for k in ["世界观", "设定", "魔法体系", "历史背景"]):
            return "worldbuilding"
        return "chapter"

    # 法学论文
    if any(k in head for k in ["摘要", "关键词", "参考文献", "作者简介"]):
        if any(k in head for k in ["实证", "数据分析", "问卷", "样本"]):
            return "paper_empirical"
        if any(k in head for k in ["比较法", "比较研究", "域外"]):
            return "paper_comparative"
        if any(k in head for k in ["案例分析", "裁判", "判决"]):
            return "paper_case"
        return "paper_thematic"

    return "paper_thematic"


def guess_domain(text: str, filename: str) -> str:
    """推测所属领域"""
    # 法学特征：有摘要/关键词/参考文献/法条
    law_indicators = 0
    if "摘要" in text[:5000]:
        law_indicators += 1
    if "关键词" in text[:5000]:
        law_indicators += 1
    if "参考文献" in text[-2000:]:
        law_indicators += 1
    if re.search(r"第\s*\d+\s*条", text):
        law_indicators += 1
    if re.search(r"[\u4e00-\u9fff]+法[第章节]", text[:1000]):
        law_indicators += 1
    # 学术论文通有词汇（在开头 3000 字中至少命中 1 个即 +1）
    academic_tone = ["本文", "理论", "制度", "数据", "治理", "规制", "规范", "体系", "建构"]
    if any(kw in text[:3000] for kw in academic_tone) and ("摘要" in text[:5000] or "关键词" in text[:5000]):
        law_indicators += 1

    # 写作特征：章节标题、对话、叙事
    writing_indicators = 0

    # 章节标记（多种格式）
    # 学术论文也有"第一章"等标记，当已检测到摘要+关键词时不加重权重
    if re.search(r"第[一二三四五六七八九十百千\d]+[章节回]", text):
        if law_indicators >= 2:
            writing_indicators += 1  # 学术论文中的章节标记仅 +1
        else:
            writing_indicators += 2  # 小说章节标记权重更高
    if re.search(r"(楔子|序章|序言|尾声|番外|后记)", text[:500]):
        writing_indicators += 2

    # 中文引号对话（"" 而非「」）
    dialogue_count = len(re.findall(r"\u201c[^\u201d]{2,}\u201d", text[:3000]))
    if dialogue_count >= 3:
        writing_indicators += 2
    elif dialogue_count >= 1:
        writing_indicators += 1

    # 「」引号对话
    if re.search(r"「[^」]+」", text[:3000]):
        writing_indicators += 1

    # 叙事特征：场景描写关键词
    narrative_patterns = [
        r"[她他]的[眼嘴角手指]}",  # 肢体描写
        r"(走进|推开|打开|转身|低头|抬头|叹了口气)",  # 动作描写
        r"(阳光|月光|灯光|雨|风|雪).{0,10}(照|洒|打|吹|落)在",  # 环境描写
    ]
    for pat in narrative_patterns:
        if re.search(pat, text[:3000]):
            writing_indicators += 1
            break

    # 文件名特征
    fn_lower = filename.lower()
    if any(k in fn_lower for k in ["chapter", "小说", "章节", "序章", "楔子", "番外", "大纲", "设定"]):
        writing_indicators += 2

    # 判断逻辑
    if law_indicators >= 3:
        return "law"  # 强法学特征
    if law_indicators >= 2 and writing_indicators < 2:
        return "law"
    if writing_indicators >= 3 and writing_indicators >= law_indicators + 2:
        return "writing"  # 需要明显的小说特征才能覆盖法学判定
    if writing_indicators > law_indicators:
        return "writing"
    return "law"  # 默认法学


# ─── 多格式文本提取 ───

def extract_text_from_pdf(filepath: Path) -> str:
    """从 PDF 提取文本"""
    doc = fitz.open(str(filepath))
    text = ""
    for page in doc:
        text += page.get_text()
    doc.close()
    return text


def extract_text_from_txt(filepath: Path) -> str:
    """从 TXT/Markdown 提取文本"""
    for encoding in ("utf-8", "gbk", "gb2312", "latin-1"):
        try:
            return filepath.read_text(encoding=encoding)
        except (UnicodeDecodeError, UnicodeError):
            continue
    return filepath.read_bytes().decode("utf-8", errors="replace")


def extract_text_from_docx(filepath: Path) -> str:
    """从 DOCX 提取文本"""
    from docx import Document
    doc = Document(str(filepath))
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    # 表格中的文本
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    parts.append(cell.text)
    return "\n".join(parts)


def extract_text_from_epub(filepath: Path) -> str:
    """从 EPUB 提取文本"""
    import ebooklib
    from ebooklib import epub
    from bs4 import BeautifulSoup

    book = epub.read_epub(str(filepath), options={"ignore_ncx": True})
    parts = []
    for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
        soup = BeautifulSoup(item.get_body_content(), "html.parser")
        text = soup.get_text(separator="\n", strip=True)
        if text:
            parts.append(text)
    return "\n\n".join(parts)


def extract_text_from_html(filepath: Path) -> str:
    """从 HTML 提取文本"""
    from bs4 import BeautifulSoup

    raw = extract_text_from_txt(filepath)
    soup = BeautifulSoup(raw, "html.parser")
    # 移除脚本和样式
    for tag in soup(["script", "style", "nav", "footer", "header"]):
        tag.decompose()
    return soup.get_text(separator="\n", strip=True)


def extract_text_from_pptx(filepath: Path) -> str:
    """从 PPTX 提取文本"""
    from pptx import Presentation
    prs = Presentation(str(filepath))
    parts = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            slide_texts.append(cell.text.strip())
        if slide_texts:
            parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_texts))
    return "\n\n".join(parts)


def extract_text_from_image(filepath: Path) -> str:
    """从图片提取文本（OCR）"""
    import pytesseract
    from PIL import Image

    img = Image.open(str(filepath))
    # 同时识别中文和英文
    text = pytesseract.image_to_string(img, lang="chi_sim+eng")
    return text.strip()


def extract_text(filepath: Path) -> str:
    """根据文件扩展名分派到对应的提取器"""
    ext = filepath.suffix.lower()
    extractors = {
        ".pdf": extract_text_from_pdf,
        ".txt": extract_text_from_txt,
        ".md": extract_text_from_txt,
        ".markdown": extract_text_from_txt,
        ".docx": extract_text_from_docx,
        ".epub": extract_text_from_epub,
        ".html": extract_text_from_html,
        ".htm": extract_text_from_html,
        ".pptx": extract_text_from_pptx,
        ".jpg": extract_text_from_image,
        ".jpeg": extract_text_from_image,
        ".png": extract_text_from_image,
        ".bmp": extract_text_from_image,
        ".tiff": extract_text_from_image,
        ".webp": extract_text_from_image,
    }
    extractor = extractors.get(ext)
    if not extractor:
        raise ValueError(f"不支持的文件格式: {ext}")
    return extractor(filepath)


def extract_archive(filepath: Path, extract_dir: Path) -> list:
    """解压压缩包，返回其中支持的文件列表"""
    import zipfile
    import tarfile

    name = filepath.name.lower()
    if name.endswith(".zip"):
        with zipfile.ZipFile(str(filepath), "r") as zf:
            zf.extractall(str(extract_dir))
    elif name.endswith((".tar.gz", ".tgz", ".tar.bz2", ".tar.xz", ".tar")):
        with tarfile.open(str(filepath), "r:*") as tf:
            tf.extractall(str(extract_dir))
    else:
        return []

    # 遍历解压后的文件，找出支持的格式（跳过嵌套压缩包）
    supported_files = []
    for f in sorted(extract_dir.rglob("*")):
        if f.is_file() and f.suffix.lower() in SUPPORTED_EXTENSIONS and not is_archive(f):
            supported_files.append(f)

    return supported_files


def process_archive(filepath: Path, kb: KnowledgeBase, db_tags: set) -> bool:
    """解压并处理压缩包中的所有支持文件"""
    import tempfile

    filename = filepath.name
    log.info(f"开始处理压缩包: {filename}")

    # 从压缩包文件名提取可能的小说名
    archive_base = filename
    for ext in [".zip", ".tar.gz", ".tgz", ".tar.bz2", ".tar.xz", ".tar"]:
        if archive_base.lower().endswith(ext):
            archive_base = archive_base[: -len(ext)]
            break

    try:
        with tempfile.TemporaryDirectory() as tmpdir:
            tmpdir = Path(tmpdir)
            files = extract_archive(filepath, tmpdir)
            log.info(f"  解压完成，发现 {len(files)} 个支持的文件")

            if not files:
                log.warning(f"  压缩包中无支持的文件: {filename}")
                return False

            success = 0
            for f in files:
                rel_path = f.relative_to(tmpdir)
                log.info(f"  处理: {rel_path}")

                # 从路径提取小说名：优先用第一级目录名，其次用压缩包名
                parts = rel_path.parts
                if len(parts) > 1:
                    book_title = parts[0]
                else:
                    book_title = archive_base

                if process_file(f, kb, db_tags, book_title=book_title):
                    success += 1

            log.info(f"  压缩包处理完成: {success}/{len(files)} 个文件成功入库")

            # 归档压缩包
            dest = ARCHIVE_DIR / f"{datetime.now().strftime('%Y%m%d_%H%M%S')}_{filename}"
            try:
                shutil.move(str(filepath), str(dest))
                log.info(f"  已归档: {dest.name}")
            except FileNotFoundError:
                log.warning(f"  文件已不存在: {filepath}")

            return success > 0
    except Exception as e:
        log.error(f"  压缩包处理失败: {filename} — {e}", exc_info=True)
        return False


def process_file(filepath: Path, kb: KnowledgeBase, db_tags: set, book_title: str = "") -> bool:
    """处理单个文件（支持多种格式）

    Args:
        filepath: 文件路径
        kb: 知识库实例
        db_tags: 数据库中的标签名集合
        book_title: 所属小说/作品名（从压缩包目录名提取），可选
    """
    filename = filepath.name

    # 压缩包走专用流程
    if is_archive(filepath):
        return process_archive(filepath, kb, db_tags)

    ext = filepath.suffix.lower()
    log.info(f"开始处理: {filename} (格式: {ext})")

    try:
        # 1. 解析文件
        log.info(f"  [1/5] 解析文件 ({ext})...")
        text = extract_text(filepath)

        if not text.strip():
            log.warning(f"  ⚠ 文本为空: {filename}")
            return False

        log.info(f"        {len(text):,} 字符")

        # 1.5 读取用户提供的元数据（.meta.json 侧车文件）
        user_meta = {}
        meta_filepath = filepath.with_suffix(filepath.suffix + ".meta.json")
        if meta_filepath.exists():
            try:
                import json as _json
                user_meta = _json.loads(meta_filepath.read_text(encoding="utf-8"))
                log.info(f"  [1.5] 读取到用户元数据: {user_meta}")
            except Exception as e:
                log.warning(f"  ⚠ 元数据文件读取失败: {e}")

        # 2. 提取元数据
        log.info(f"  [2/5] 提取元数据...")
        title = user_meta.get("title", "") or extract_title(text, filename)
        author = user_meta.get("author", "") or extract_author(text, filename)
        # 用户指定的 domain 优先于自动检测
        if user_meta.get("domain"):
            domain = user_meta["domain"]
        else:
            domain = guess_domain(text, filename)
        summary = extract_summary(text)
        keywords = extract_keywords(text)
        year = extract_year(text)
        doc_type = guess_doc_type(text, filename, domain)

        log.info(f"       标题: {title[:60]}")
        log.info(f"       作者: {author}")
        log.info(f"       领域: {domain} | 类型: {doc_type}")
        log.info(f"       关键词: {keywords}")

        # 3. 自动标签 + 用户标签合并
        log.info(f"  [3/5] 自动打标签...")
        tags = auto_tag_enhanced(text, title, keywords, db_tags, domain=domain)
        # 合并用户手动指定的标签（逗号分隔）
        if user_meta.get("tags"):
            user_tags = [t.strip() for t in user_meta["tags"].split(",") if t.strip()]
            # 用户标签优先，去重后放在前面
            existing_set = set(tags)
            for ut in user_tags:
                if ut not in existing_set:
                    tags.insert(0, ut)
                    existing_set.add(ut)
            log.info(f"       合并用户标签: {user_tags}")
        log.info(f"       最终标签: {tags}")

        # 4. 入库
        log.info(f"  [4/5] 录入数据库...")
        published_at = f"{year}-01-01" if year else None
        meta = {
            "filename": filename,
            "format": ext,
            "char_count": len(text),
            "keywords": keywords,
            "year": year,
            "ingested_by": "kb_watcher",
        }
        if book_title:
            meta["book_title"] = book_title
            log.info(f"       所属作品: {book_title}")
        result = kb.ingest(
            title=title,
            domain=domain,
            doc_type=doc_type,
            content=text,
            source="自动监控入库",
            source_url=None,
            author=author,
            published_at=published_at,
            summary=summary,
            metadata=meta,
            tags=tags,
        )
        doc_id = result.get("id", "?")
        log.info(f"  ✅ 入库成功: {doc_id[:8]}... ({title[:40]})")

        # 4.5 抽取关键概念（KeyBERT + 知识复利）
        log.info(f"  [4.5/5] 抽取关键概念...")
        try:
            concepts = kb.extract_concepts(doc_id, top_n=10)
            if concepts:
                top_names = [c["name"] for c in concepts[:5]]
                new_count = sum(1 for c in concepts if c.get("is_new"))
                log.info(f"       抽取到 {len(concepts)} 个概念: {', '.join(top_names)}{'...' if len(concepts) > 5 else ''}")
            else:
                log.info(f"       ⚠ 未抽取到概念（文本可能过短或模型未就绪）")
        except Exception as e:
            log.warning(f"  ⚠ 概念抽取失败（不影响入库）: {e}")

        # 4.6 生成结构化摘要（P1）
        log.info(f"  [4.6/5] 生成结构化摘要...")
        try:
            summary = kb.generate_summary(doc_id)
            if summary and not summary.get("error"):
                log.info(f"       摘要已生成: {summary.get('core_argument', '')[:60]}...")
            else:
                log.info(f"       ⚠ 摘要生成跳过")
        except Exception as e:
            log.warning(f"  ⚠ 摘要生成失败（不影响入库）: {e}")

        # 4.7 记录操作日志（P0）
        try:
            kb.log_operation(
                operation_type="ingest",
                entity_type="document",
                entity_id=doc_id,
                entity_title=title,
                details={
                    "domain": domain,
                    "doc_type": doc_type,
                    "char_count": len(text),
                    "tags": tags,
                    "concepts_extracted": len(concepts) if 'concepts' in dir() else 0,
                },
                operator="watcher",
            )
        except Exception:
            pass

        # 5. 归档
        log.info(f"  [5/5] 归档到 {ARCHIVE_DIR}/")
        dest = ARCHIVE_DIR / f"{datetime.now().strftime('%Y%m%d_%H%M%S')}_{filename}"
        try:
            shutil.move(str(filepath), str(dest))
            log.info(f"  📦 已归档: {dest.name}")
            # 同时归档元数据侧车文件（如果存在）
            if meta_filepath.exists():
                try:
                    meta_dest = ARCHIVE_DIR / f"{datetime.now().strftime('%Y%m%d_%H%M%S')}_{meta_filepath.name}"
                    shutil.move(str(meta_filepath), str(meta_dest))
                    log.info(f"  📦 已归档元数据: {meta_dest.name}")
                except FileNotFoundError:
                    pass
        except FileNotFoundError:
            log.warning(f"  ⚠ 文件已不存在（可能已被其他进程处理）: {filepath}")

        return True

    except Exception as e:
        log.error(f"  ❌ 处理失败: {filename} — {e}", exc_info=True)
        return False


def scan_existing(observer):
    """扫描 WATCH_DIR 中已存在的文件"""
    for filepath in sorted(WATCH_DIR.iterdir()):
        if filepath.suffix.lower() in SUPPORTED_EXTENSIONS:
            # 触发 handler 的 on_created
            observer.schedule(
                observer._handlers[0], str(WATCH_DIR), recursive=False
            )
            yield filepath


class FileHandler(FileSystemEventHandler):
    """文件监控处理器"""

    def __init__(self, kb, db_tags):
        self.kb = kb
        self.db_tags = db_tags
        self.processing = set()

    def on_created(self, event):
        if event.is_directory:
            return
        filepath = Path(event.src_path)
        if filepath.suffix.lower() not in SUPPORTED_EXTENSIONS and not is_archive(filepath):
            return
        if filepath.name in self.processing:
            return

        # 等待文件写入完成（避免拿到不完整的文件）
        self.processing.add(filepath.name)
        time.sleep(2)  # 给文件系统一点时间完成写入

        # 如果文件在子目录中，用父目录名作为 book_title
        try:
            rel = filepath.relative_to(WATCH_DIR)
            parts = rel.parts
            book_title = parts[0] if len(parts) > 1 else ""
        except ValueError:
            book_title = ""

        try:
            process_file(filepath, self.kb, self.db_tags, book_title=book_title)
        finally:
            self.processing.discard(filepath.name)


def main():
    log.info("=" * 60)
    log.info("知识库文件监控启动")
    log.info(f"监控目录: {WATCH_DIR}")
    log.info(f"归档目录: {ARCHIVE_DIR}")
    log.info("=" * 60)

    # 确保目录存在
    WATCH_DIR.mkdir(parents=True, exist_ok=True)
    ARCHIVE_DIR.mkdir(parents=True, exist_ok=True)

    # 连接数据库
    kb = KnowledgeBase()
    with kb.conn.cursor() as cur:
        cur.execute("SELECT name FROM tags")
        db_tags = {r[0] for r in cur.fetchall()}
    log.info(f"数据库就绪，标签 {len(db_tags)} 个（法学+写作）")

    # 处理已有文件
    log.info("扫描已有文件...")
    existing = []
    for f in sorted(WATCH_DIR.rglob("*")):
        if f.is_file() and (f.suffix.lower() in SUPPORTED_EXTENSIONS or is_archive(f)):
            existing.append(f)
    for fp in sorted(existing):
        log.info(f"发现已有文件: {fp.name}")
        # 提取 book_title（子目录名）
        try:
            rel = fp.relative_to(WATCH_DIR)
            parts = rel.parts
            book_title = parts[0] if len(parts) > 1 else ""
        except ValueError:
            book_title = ""
        process_file(fp, kb, db_tags, book_title=book_title)

    # 启动监控
    handler = FileHandler(kb, db_tags)
    observer = Observer()
    observer.schedule(handler, str(WATCH_DIR), recursive=True)
    observer.start()

    log.info("👀 文件监控已启动，等待新文件...")
    log.info(f"支持格式: {', '.join(sorted(SUPPORTED_EXTENSIONS))}, 压缩包: {', '.join(sorted(ARCHIVE_EXTENSIONS))}")
    log.info("按 Ctrl+C 停止")

    def shutdown(signum, frame):
        log.info("\n收到停止信号，关闭...")
        observer.stop()
        observer.join()
        log.info("已停止")
        sys.exit(0)

    signal.signal(signal.SIGINT, shutdown)
    signal.signal(signal.SIGTERM, shutdown)

    try:
        while True:
            time.sleep(1)
    except KeyboardInterrupt:
        shutdown(None, None)


if __name__ == "__main__":
    main()
